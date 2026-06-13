#!/usr/bin/env python3
"""Generate a polished LastMinute ITR investor pitch deck (PowerPoint)."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

OUTPUT = Path(__file__).resolve().parent / "LastMinute_ITR_Pitch_Deck.pptx"

# Brand palette
PRIMARY = RGBColor(37, 99, 235)
PRIMARY_DARK = RGBColor(30, 58, 138)
ACCENT = RGBColor(96, 165, 250)
TEXT = RGBColor(15, 23, 42)
MUTED = RGBColor(100, 116, 139)
WHITE = RGBColor(255, 255, 255)
LIGHT = RGBColor(248, 250, 252)
BORDER = RGBColor(226, 232, 240)
GREEN = RGBColor(22, 163, 74)
ORANGE = RGBColor(234, 88, 12)


def blank_slide(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def solid_bg(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def rect(slide, left, top, width, height, color, *, radius=False):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def textbox(slide, left, top, width, height, text, *, size=18, bold=False, color=TEXT, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return box


def bullets(slide, items, left, top, width, height, *, size=16, color=TEXT, spacing=8):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(spacing)
        p.bullet = True
    return box


def slide_header(slide, prs, title, subtitle=None, dark=False):
    if dark:
        solid_bg(slide, PRIMARY_DARK)
        title_color, sub_color = WHITE, ACCENT
        accent = ACCENT
    else:
        solid_bg(slide, WHITE)
        rect(slide, 0, 0, prs.slide_width, Inches(0.14), PRIMARY)
        rect(slide, Inches(0.55), Inches(0.55), Inches(0.08), Inches(0.95), PRIMARY)
        title_color, sub_color = TEXT, MUTED
        accent = PRIMARY

    textbox(slide, Inches(0.8), Inches(0.45), Inches(8.5), Inches(0.7), title, size=34, bold=True, color=title_color)
    if subtitle:
        textbox(slide, Inches(0.8), Inches(1.15), Inches(8.5), Inches(0.55), subtitle, size=16, color=sub_color)
    textbox(slide, Inches(0.55), Inches(7.05), Inches(8.8), Inches(0.3), "LastMinute ITR · Confidential · AY 2026-27", size=9, color=MUTED if not dark else RGBColor(148, 163, 184))
    return accent


def stat_cards(slide, stats, top=Inches(2.0)):
    card_w = Inches(2.55)
    gap = Inches(0.22)
    x = Inches(0.7)
    for value, label, detail in stats:
        rect(slide, x, top, card_w, Inches(1.6), LIGHT, radius=True)
        rect(slide, x, top, card_w, Inches(0.08), PRIMARY, radius=True)
        textbox(slide, x + Inches(0.18), top + Inches(0.22), card_w - Inches(0.36), Inches(0.45), value, size=28, bold=True, color=PRIMARY)
        textbox(slide, x + Inches(0.18), top + Inches(0.72), card_w - Inches(0.36), Inches(0.3), label, size=12, bold=True)
        textbox(slide, x + Inches(0.18), top + Inches(1.02), card_w - Inches(0.36), Inches(0.35), detail, size=10, color=MUTED)
        x += card_w + gap


def slide_cover(prs):
    slide = blank_slide(prs)
    solid_bg(slide, PRIMARY_DARK)
    rect(slide, Inches(6.5), Inches(-0.5), Inches(4), Inches(4), RGBColor(29, 78, 216))
    rect(slide, Inches(7.8), Inches(4.8), Inches(3), Inches(3), RGBColor(37, 99, 235))

    circle = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(0.85), Inches(1.35), Inches(0.9), Inches(0.9))
    circle.fill.solid()
    circle.fill.fore_color.rgb = WHITE
    circle.line.fill.background()
    textbox(slide, Inches(0.85), Inches(1.55), Inches(0.9), Inches(0.5), "LM", size=22, bold=True, color=PRIMARY, align=PP_ALIGN.CENTER)

    textbox(slide, Inches(0.85), Inches(2.45), Inches(8), Inches(0.9), "LastMinute ITR", size=46, bold=True, color=WHITE)
    textbox(slide, Inches(0.85), Inches(3.35), Inches(7.5), Inches(0.6), "Your friendly AI CA for last-minute ITR filing", size=22, color=ACCENT)
    textbox(slide, Inches(0.85), Inches(4.05), Inches(7.2), Inches(0.5), "Prep with AI. You file on incometax.gov.in yourself.", size=18, color=WHITE)

    pill = rect(slide, Inches(0.85), Inches(4.85), Inches(3.2), Inches(0.55), RGBColor(30, 58, 138), radius=True)
    textbox(slide, Inches(1.05), Inches(4.97), Inches(2.8), Inches(0.35), "Investor Pitch Deck · June 2026", size=12, color=ACCENT)

    textbox(slide, Inches(0.85), Inches(6.2), Inches(8), Inches(0.4), "Lawful tax saving · Proof-based claims · DPDP compliant", size=11, color=RGBColor(148, 163, 184))


def slide_problem(prs):
    slide = blank_slide(prs)
    slide_header(slide, prs, "The Problem", "7+ crore Indians file ITR every year — most dread it")
    stat_cards(slide, [
        ("7.28 Cr", "ITRs filed", "AY 2024-25 official data"),
        ("3.34 Cr", "ITR-1 filers", "Primary salaried wedge"),
        ("30 days", "E-verify window", "Miss it → invalid return"),
    ])
    bullets(slide, [
        "Wrong form anxiety — ITR-1 vs 2 vs 4 confusion after rule changes",
        "AIS / Form 26AS mismatches cause refund delays and notices",
        "Old vs new regime choice is opaque; users leave money on the table",
        "Competitors over-promise auto-filing; trust breaks when data is wrong",
        "CAs are expensive; DIY portals use jargon normal users can't follow",
    ], Inches(0.9), Inches(3.9), Inches(8.3), Inches(2.8), size=15)


def slide_solution(prs):
    slide = blank_slide(prs)
    slide_header(slide, prs, "Our Solution", "AI-assisted ITR prep with companion-mode filing on the government portal")
    bullets(slide, [
        "Import-first funnel: Form 16 → parse → regime compare → paywall → portal guide",
        "Deterministic Python tax engine (ITR-1–7 rules) before any LLM layer",
        "Mismatch detection across AIS, 26AS, and Form 16 before you file",
        "Lawful deduction suggestions — proof-based, never fake claims",
        "Companion mode: screen-by-screen walkthrough on incometax.gov.in",
        "User always files, submits, and e-verifies themselves — we never auto-submit",
    ], Inches(0.9), Inches(1.95), Inches(8.3), Inches(4.8), size=16)


def slide_how(prs):
    slide = blank_slide(prs)
    slide_header(slide, prs, "How It Works", "4-step journey from panic to filed return")
    steps = [
        ("1", "Upload", "Form 16, AIS,\nbank statements"),
        ("2", "Review", "Regime compare,\nmismatch fixes"),
        ("3", "Unlock", "Pay once for\nportal guide"),
        ("4", "File", "Copy to incometax.gov.in\n& e-verify"),
    ]
    x = Inches(0.65)
    w = Inches(2.15)
    for num, title, desc in steps:
        rect(slide, x, Inches(2.05), w, Inches(2.35), LIGHT, radius=True)
        circle = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, x + Inches(0.82), Inches(2.25), Inches(0.5), Inches(0.5))
        circle.fill.solid()
        circle.fill.fore_color.rgb = PRIMARY
        circle.line.fill.background()
        textbox(slide, x + Inches(0.82), Inches(2.37), Inches(0.5), Inches(0.35), num, size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        textbox(slide, x + Inches(0.15), Inches(2.95), w - Inches(0.3), Inches(0.4), title, size=18, bold=True, align=PP_ALIGN.CENTER)
        textbox(slide, x + Inches(0.15), Inches(3.4), w - Inches(0.3), Inches(0.8), desc, size=12, color=MUTED, align=PP_ALIGN.CENTER)
        if num != "4":
            arrow = textbox(slide, x + w - Inches(0.05), Inches(2.95), Inches(0.3), Inches(0.4), "→", size=20, bold=True, color=PRIMARY)
        x += w + Inches(0.12)

    rect(slide, Inches(0.7), Inches(4.75), Inches(8.5), Inches(1.35), RGBColor(239, 246, 255), radius=True)
    textbox(slide, Inches(0.95), Inches(4.95), Inches(8.0), Inches(1.0),
            "Companion Mode = trust moat\nWe prepare numbers and guide screen-by-screen. User retains control. No ERI dependency at launch.",
            size=14, color=PRIMARY_DARK)


def slide_market(prs):
    slide = blank_slide(prs)
    slide_header(slide, prs, "Market Opportunity", "Large, recurring, deadline-driven demand")
    stat_cards(slide, [
        ("₹3–10 Cr", "Year 2–3", "3L–5L users @ 12–18% paid"),
        ("₹80L–1.5Cr", "Good Year 1", "1L users @ 10–15% paid"),
        ("₹15–50L", "Conservative Y1", "25k–50k users @ 6–10% paid"),
    ], top=Inches(1.95))
    bullets(slide, [
        "Wedge: salaried + pension + interest income (ITR-1 / ITR-2)",
        "Seasonal spike Jun–Jul filing season → predictable CAC windows",
        "Expand to freelancers (44ADA), capital gains, notices, annual plans",
        "B2B: employer partnerships, payroll integrations, ERI/API path later",
    ], Inches(0.9), Inches(3.85), Inches(8.3), Inches(2.8), size=15)


def slide_business(prs):
    slide = blank_slide(prs)
    slide_header(slide, prs, "Business Model", "Freemium funnel → paid portal companion unlock")
    rows = [
        ("Free", "₹0", "Tax estimate, ITR form selector, filing checklist"),
        ("DIY", "₹499", "Form 16 import, wizard, pre-submit checks, portal guide"),
        ("AI Smart ★", "₹349 launch", "Mismatch engine, regime optimizer, lawful deductions"),
        ("CA Review", "₹2,499", "Expert sign-off, notice-risk review (waitlist)"),
    ]
    y = Inches(2.0)
    for i, (plan, price, desc) in enumerate(rows):
        bg = RGBColor(239, 246, 255) if "AI Smart" in plan else LIGHT
        rect(slide, Inches(0.7), y, Inches(8.5), Inches(0.85), bg, radius=True)
        textbox(slide, Inches(0.95), y + Inches(0.18), Inches(1.5), Inches(0.4), plan, size=15, bold=True)
        textbox(slide, Inches(2.5), y + Inches(0.18), Inches(1.2), Inches(0.4), price, size=15, bold=True, color=PRIMARY)
        textbox(slide, Inches(3.8), y + Inches(0.18), Inches(5.2), Inches(0.5), desc, size=13, color=MUTED)
        y += Inches(0.95)
    textbox(slide, Inches(0.7), Inches(5.9), Inches(8.5), Inches(0.5),
            "Unit economics: ₹349–899 ARPU · Razorpay live · Upsell to complex cases ₹2,999–9,999+", size=12, color=MUTED)


def slide_competition(prs):
    slide = blank_slide(prs)
    slide_header(slide, prs, "Competitive Landscape", "Win on trust + simplicity, not AI hype")
    headers = ["", "ClearTax / Quicko", "Traditional CA", "LastMinute ITR"]
    rows = [
        ["Price", "₹799–999+", "₹2,000–5,000", "₹349–899"],
        ["Filing model", "Platform-led", "Human-led", "Companion (user files)"],
        ["Mismatch checks", "Varies", "Manual", "Engine-first, pre-file"],
        ["Trust", "Auto-file implied", "High trust, slow", "Honest, proof-based"],
        ["Last-minute UX", "Cluttered", "Unavailable", "Built for deadline panic"],
    ]
    col_x = [0.7, 2.2, 4.0, 6.0]
    col_w = [1.35, 1.65, 1.85, 2.3]
    y = 1.95
    for col, h in zip(col_x, headers):
        textbox(slide, Inches(col), Inches(y), Inches(col_w[col_x.index(col)]), Inches(0.35), h, size=11, bold=True, color=PRIMARY_DARK)
    for r, row in enumerate(rows):
        ry = y + 0.42 + r * 0.72
        for col, cell in zip(col_x, row):
            color = GREEN if col == 6.0 and r > 0 else TEXT
            textbox(slide, Inches(col), Inches(ry), Inches(col_w[col_x.index(col)]), Inches(0.6), cell, size=12, bold=(col == 0), color=color)


def slide_traction(prs):
    slide = blank_slide(prs)
    slide_header(slide, prs, "Traction & Product Status", "MVP live — iterating toward launch scale")
    bullets(slide, [
        "✓ Live on Vercel with Next.js app + Python L1 tax engine",
        "✓ Full funnel: import → regime → deductions → paywall → companion",
        "✓ Razorpay payments integrated; companion unlock on payment",
        "✓ SEO content: learn guides, glossary, blogs for last-minute ITR keywords",
        "✓ Analytics + E2E test suite in place",
        "→ In progress: Apple/Quicko UI polish, funnel simplification, CA waitlist",
        "→ Next: closed beta (100–300 users), CA audit loop, performance marketing",
    ], Inches(0.9), Inches(1.95), Inches(8.3), Inches(4.5), size=15)
    rect(slide, Inches(0.7), Inches(5.5), Inches(8.5), Inches(0.7), LIGHT, radius=True)
    textbox(slide, Inches(0.95), Inches(5.65), Inches(8.0), Inches(0.4),
            "Live demo: lastminute-reb94pi2k-nikhil-anand-s-projects12.vercel.app", size=13, color=PRIMARY)


def slide_tech(prs):
    slide = blank_slide(prs)
    slide_header(slide, prs, "Technology & Moat", "Engine-first architecture builds defensible accuracy")
    bullets(slide, [
        "Python L1 engine from official ITR-1–7 forms — deterministic math",
        "LLM layer only for explanations & Q&A, not tax computation",
        "Companion portal guide mapped to incometax.gov.in screen flow",
        "DPDP-compliant data handling, consent, encryption, audit logs",
        "Mismatch reconciliation engine (AIS / 26AS / Form 16) = core product",
        "Future: ERI/API filing path when certified (Income Tax Dept Type 2 ERI)",
    ], Inches(0.9), Inches(1.95), Inches(8.3), Inches(4.5), size=15)
    textbox(slide, Inches(0.7), Inches(5.4), Inches(8.5), Inches(0.5),
            "Stack: Next.js · Python tax engine · Razorpay · Vercel", size=12, color=MUTED)


def slide_gtm(prs):
    slide = blank_slide(prs)
    slide_header(slide, prs, "Go-To-Market", "Own the last-minute ITR moment")
    bullets(slide, [
        "SEO: deadline pages, regime compare, Form 16 guides, glossary",
        "Performance ads in Jun–Jul filing season with ₹349 launch offer",
        "Content: Reddit/YouTube — mismatch fear, regime confusion",
        "Referral: share regime savings card, WhatsApp-friendly checklist",
        "Partnerships: payroll SaaS, fintech apps, employer HR programs",
        "Retention: annual plan (reminders, planning, deduction nudges)",
    ], Inches(0.9), Inches(1.95), Inches(8.3), Inches(4.8), size=15)


def slide_roadmap(prs):
    slide = blank_slide(prs)
    slide_header(slide, prs, "Roadmap", "Narrow wedge → trust → expand")
    phases = [
        ("Now", "Salaried ITR-1/2, companion mode, AI Smart launch"),
        ("Q3 2026", "Closed beta, CA review waitlist, funnel optimization"),
        ("Q4 2026", "Freelancer 44ADA, senior citizen flows, annual plan"),
        ("2027", "Capital gains, notices, ERI/API path, B2B employer channel"),
    ]
    y = 2.0
    for i, (phase, desc) in enumerate(phases):
        color = PRIMARY if i == 0 else LIGHT
        text_color = WHITE if i == 0 else PRIMARY_DARK
        rect(slide, Inches(0.7), Inches(y), Inches(1.4), Inches(0.55), color, radius=True)
        textbox(slide, Inches(0.7), Inches(y + 0.1), Inches(1.4), Inches(0.35), phase, size=13, bold=True, color=text_color, align=PP_ALIGN.CENTER)
        textbox(slide, Inches(2.3), Inches(y + 0.1), Inches(6.8), Inches(0.45), desc, size=16)
        y += 1.05


def slide_ask(prs):
    slide = blank_slide(prs)
    slide_header(slide, prs, "The Ask", dark=True)
    bullets(slide, [
        "Raise seed capital to scale Jul 2026 filing season acquisition",
        "Hire: growth marketer, CA network lead, senior full-stack engineer",
        "12-month target: 50k–100k users, 8–12% paid conversion",
        "Use of funds: 40% growth · 30% product · 20% CA ops · 10% compliance",
    ], Inches(1.0), Inches(2.2), Inches(8.0), Inches(3.2), size=20, color=WHITE, spacing=14)
    textbox(slide, Inches(0.85), Inches(5.5), Inches(8.2), Inches(0.6),
            "Let's make ITR filing feel like having a personal CA — without the CA bill.",
            size=16, color=ACCENT)
    textbox(slide, Inches(0.85), Inches(6.2), Inches(8.2), Inches(0.4),
            "contact@lastminuteitr.com · lastminuteitr.com", size=13, color=RGBColor(148, 163, 184))


def build_deck() -> Path:
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    slide_cover(prs)
    slide_problem(prs)
    slide_solution(prs)
    slide_how(prs)
    slide_market(prs)
    slide_business(prs)
    slide_competition(prs)
    slide_traction(prs)
    slide_tech(prs)
    slide_gtm(prs)
    slide_roadmap(prs)
    slide_ask(prs)

    prs.save(OUTPUT)
    return OUTPUT


if __name__ == "__main__":
    print(f"Created: {build_deck()}")
